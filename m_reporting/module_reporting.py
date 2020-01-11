from pptx import Presentation
from pptx.util import Inches, Pt

import email, smtplib, ssl, getpass

from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

#----------ppt
def create_powerpoint(file_name):
    ppt = Presentation()
    save_ppt(ppt, file_name)


def save_ppt(ppt, file_name):
    ppt.save(file_name)


def open_powerpoint(file_name):
    ppt = Presentation(file_name)
    return ppt


def create_slide(ppt, layout):
    return ppt.slides.add_slide(layout)


def add_basic_layouts(ppt):
    ppt_layout = ppt.slide_layouts[0]
    ppt_layout1 = ppt.slide_layouts[1]
    ppt_layout2 = ppt.slide_layouts[2]
    ppt_layout3 = ppt.slide_layouts[3]
    ppt_layout4 = ppt.slide_layouts[4]
    ppt_layout5 = ppt.slide_layouts[5]
    ppt_layout6 = ppt.slide_layouts[6]
    ppt_layout7 = ppt.slide_layouts[7]
    ppt_layout8 = ppt.slide_layouts[8]
    ppt_layout9 = ppt.slide_layouts[9]
    ppt_layout10 = ppt.slide_layouts[10]

def add_blank_layout(ppt):
    ppt_layout = ppt.slide_layouts[6]
    current_slide = create_slide(ppt, ppt_layout)
    return current_slide

def add_text(slide, text, bold, font_size, pos_left, pos_right):
    area_left = Inches(pos_left)
    area_right = Inches(pos_right)
    area_height_width = Inches(1)
    text_area = slide.shapes.add_textbox(area_left, area_right, area_height_width, area_height_width)
    text_frame = text_area.text_frame

    current_paragraph = text_frame.add_paragraph()
    current_paragraph.text = text

    if bold == True:
        current_paragraph.font.bold = True

    current_paragraph.font.size = Pt(font_size)


def add_image(slide, image, pos_left, pos_top, width, *height):
    area_left = Inches(pos_left)
    area_top = Inches(pos_top)
    area_width = Inches(width)

    if len(height) > 0:
        area_height = Inches(height[0])
        slide.shapes.add_picture(image, area_left, area_top, area_width, area_height)
    else:
        slide.shapes.add_picture(image, area_left, area_top, area_width)


def add_table(slide, nb_columns, nb_rows, pos_left, pos_top, width, *height):
    area_left = Inches(pos_left)
    area_top = Inches(pos_top)
    area_width = Inches(width)

    if len(height) > 0:
        area_height = Inches(height[0])
    else:
        area_height = Inches(width)

    current_table = slide.shapes.add_table(nb_rows, nb_columns, area_left, area_top, area_width, area_height)

    current_table.table.cell(0, 0).text = "Column 1"
    current_table.table.cell(0, 1).text = "Column 2"
    current_table.table.cell(1, 0).text = "Row 1"
    current_table.table.cell(1, 1).text = "Row 2"


def create_ppt(file_name):
    print('7.- Creating the PowerPoint')
    file_name = './data/results/' + file_name + '.pptx'
    create_powerpoint(f'{file_name}')
    ppt = open_powerpoint(f'{file_name}')

    add_basic_layouts(ppt)


    # slide 0
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Forbes List Analysis", True, 45, 1, 4)
    add_text(current_slide, "2018", True, 30, 1, 5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 3.75, 0.4, 2.8)


    # slide 1a
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Position Vs. Worth", True, 35, 1, 4)
    add_text(current_slide, "Position as rank position", False, 20, 1, 5)
    add_text(current_slide, "Worth as billions of United State dollars", False, 20, 1, 5.3)
    add_text(current_slide, "Period: 2018", False, 20, 1, 5.6)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 1b
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Position Vs. Worth", True, 25, 3, 0.4)
    add_image(current_slide, "./data/processed/images/position_worthbusd.png", 0.7, 1.1, 7.5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)



    # slide 2a
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Region Vs. Countries", True, 35, 1, 4)
    add_text(current_slide, "Region as world regions", False, 20, 1, 5)
    add_text(current_slide, "Countries as number of countries belonging to the same world region", False, 20, 1, 5.3)
    add_text(current_slide, "Period: 2018", False, 20, 1, 5.6)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 2b
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Region Vs. Countries", True, 25, 3, 0.4)
    add_image(current_slide, "./data/processed/images/region_numcountries.png", 0.7, 2.1, 7.5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)

    # slide 3a
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Sector Vs. Area", True, 35, 1, 4)
    add_text(current_slide, "Sector as market families", False, 20, 1, 5)
    add_text(current_slide, "Area in squared milles", False, 20, 1, 5.3)
    add_text(current_slide, "Period: 2018", False, 20, 1, 5.6)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 3b
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Sector Vs. Area", True, 25, 3, 0.4)
    add_image(current_slide, "./data/processed/images/sector_area.png", 0.7, 2.1, 7.5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 4a
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Sector Vs. Corporations", True, 35, 1, 4)
    add_text(current_slide, "Sector as market families", False, 20, 1, 5)
    add_text(current_slide, "Corporations as number of corporations belonging to the same sector", False, 20, 1, 5.3)
    add_text(current_slide, "Period: 2018", False, 20, 1, 5.6)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 4b
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Sector Vs. Corporations", True, 25, 3, 0.4)
    add_image(current_slide, "./data/processed/images/sector_numcorps.png", 0.7, 2.1, 7.5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)



    # slide 5a
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Sector Vs. Countries", True, 35, 1, 4)
    add_text(current_slide, "Sector as market families", False, 20, 1, 5)
    add_text(current_slide, "Countries as number of countries operating the same sector", False, 20, 1, 5.3)
    add_text(current_slide, "Period: 2018", False, 20, 1, 5.6)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 5b
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Sector Vs. Countries", True, 25, 3, 0.4)
    add_image(current_slide, "./data/processed/images/sector_numcountries.png", 1, 2.1, 7.5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 6a
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Worth Vs. GDP", True, 35, 1, 4)
    add_text(current_slide, "Worth as billions of United State dollars", False, 20, 1, 5)
    add_text(current_slide, "GDP as Gross Domestic Product in Billions of United States dollar", False, 20, 1, 5.3)
    add_text(current_slide, "Period: 2018", False, 20, 1, 5.6)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 6b
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Worth Vs. GDP", True, 25, 3, 0.4)
    add_image(current_slide, "./data/processed/images/worthbusd_GDP.png", 0.7, 1.1, 7.5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 7a
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Worth Vs. Phones", True, 35, 1, 4)
    add_text(current_slide, "Worth as billions of United State dollars", False, 20, 1, 5)
    add_text(current_slide, "Phones as number of phones per 1000 citicens", False, 20, 1, 5.3)
    add_text(current_slide, "Period: 2018", False, 20, 1, 5.6)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)


    # slide 7b
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "Worth Vs. Phones", True, 25, 3, 0.4)
    add_image(current_slide, "./data/processed/images/worthbusd_phones.png", 0.7, 1.1, 7.5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 7.75, 0.4, 1.6)

    # slide 8
    current_slide = add_blank_layout(ppt)
    add_text(current_slide, "THANKS FOR YOUR ATTENTION", True, 45, 1, 4)
    add_text(current_slide, "HACK-REPORT", True, 20, 4, 5)
    add_image(current_slide, "./data/raw/ironhack_logo.png", 3.75, 0.4, 2.8)

    save_ppt(ppt, file_name)

#----------email


# Create a multipart message and set headers

def send_email(file_name):
    send_question = str(input('\nDo you want to send your PPT report by email? (y/n): '))
    if send_question == 'n':
        pass

    else:
        print('\n8.- Send an email? Ok, but...')
        print('...we need some more information:\n')
        sender_email = 'herreradelduque0@gmail.com'
        receiver_email = input('Please enter RECEIVER email: ')#'herreradelduque0@gmail.com'
        subject ='[New Report] Forbes Data Analysis (2018)'
        receivers_email = 'herreradelduque0@gmail.com'


        message = MIMEMultipart()
        message["From"] = sender_email
        message["To"] = receiver_email
        message["Subject"] = subject
        message["Bcc"] = receivers_email  # Recommended for mass emails

        # Add body to email

        body = 'Dear,\n\n Please find attached the new Forbes data analysis (2018). \n\n The attached file contains the main charts done in yearly basis.\n\n' \
               'We will be available for you in case of doubts with the report. \n\n Best Regards. \n Víctor Galán'

        message.attach(MIMEText(body, "plain"))

        filename  = './data/results/' + file_name + '.pptx'  # In same directory as script

        # Open PDF file in binary mode
        with open(filename, "rb") as attachment:
            # Add file as application/octet-stream
            # Email client can usually download this automatically as attachment
            part = MIMEBase("application", "octet-stream")
            part.set_payload(attachment.read())

        # Encode file in ASCII characters to send by email
        encoders.encode_base64(part)

        # Add header as key/value pair to attachment part
        part.add_header(
            "Content-Disposition",
            f"attachment; filename= {filename}",
        )

        # Add attachment to message and convert message to string
        message.attach(part)
        text = message.as_string()

        # Log in to server using secure context and send email

        password = getpass.getpass(prompt='What is your email password? ')

        context = ssl.create_default_context()
        with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=context) as server:
            server.login(sender_email, password)
            server.sendmail(sender_email, receiver_email, text)
        print('\nWe have sent the reporting email to the receiver.')